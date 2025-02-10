import pandas as pd
import numpy as np
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import datetime
import comtypes.client
import os
import io


def convert_ppt_to_pdf(input_file, output_file):
    #Inicialização COM
    comtypes.CoInitialize()
    # Verificar se o caminho é absoluto
    input_file = os.path.abspath(input_file)
    output_file = os.path.abspath(output_file)

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.DisplayAlerts = 0  # Desativar alertas
    presentation = powerpoint.Presentations.Open(input_file, WithWindow=False)
    presentation.SaveAs(output_file, 32)  # 32 is the format type for PDF
    presentation.Close()
    powerpoint.Quit()
    #Finalizar COM
    comtypes.CoUninitialize()


# Função para formatar texto
def formatar_texto(run, texto, fonte, tamanho):
    run.text = texto
    font = run.font
    font.name = fonte
    font.size = Pt(tamanho)
    run.alignment = PP_ALIGN.CENTER  # Alinha o texto no centro


# Função para formatar texto
def formatar_texto_proposta(run, texto, fonte, tamanho):
    run.text = texto
    font = run.font
    font.name = fonte
    font.size = Pt(tamanho)


def atualizar_texto_shape_proposta(shape, texto1, fonte1, tamanho1, texto2=None, fonte2=None, tamanho2=None):
    # Define fonte e tamanho padrão para texto2 se não forem fornecidos
    if fonte2 is None:
        fonte2 = fonte1
    if tamanho2 is None:
        tamanho2 = tamanho1

    # Obtém o TextFrame do shape
    text_frame = shape.text_frame
    text_frame.word_wrap = True  # Ativa quebra de linha automática

    # Atualiza o texto1 mantendo a formatação existente
    if text_frame.paragraphs:
        p = text_frame.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            formatar_texto_proposta(run, texto1, fonte1, tamanho1)
        else:
            p.text = texto1
            p.font.name = fonte1
            p.font.size = Pt(tamanho1)

    # Atualiza o texto2 mantendo a formatação existente, se fornecido
    if texto2:
        if len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[1]
            if p.runs:
                run = p.runs[0]
                formatar_texto_proposta(run, texto2, fonte2, tamanho2)
            else:
                p.text = texto2
                p.font.name = fonte2
                p.font.size = Pt(tamanho2)

        else:
            p = text_frame.add_paragraph()
            p.text = texto2
            p.font.name = fonte2
            p.font.size = Pt(tamanho2)


def atualizar_texto_shape(shape, texto1, fonte1, tamanho1, texto2=None, fonte2=None, tamanho2=None):
    # Define fonte e tamanho padrão para texto2 se não forem fornecidos
    if fonte2 is None:
        fonte2 = fonte1
    if tamanho2 is None:
        tamanho2 = tamanho1

    # Obtém o TextFrame do shape
    text_frame = shape.text_frame
    text_frame.word_wrap = True  # Ativa quebra de linha automática
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Alinha o parágrafo no centro

    # Atualiza o texto1 mantendo a formatação existente
    if text_frame.paragraphs:
        p = text_frame.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            formatar_texto(run, texto1, fonte1, tamanho1)
        else:
            p.text = texto1
            p.font.name = fonte1
            p.font.size = Pt(tamanho1)
            p.alignment = PP_ALIGN.CENTER  # Alinha o texto no centro

    # Atualiza o texto2 mantendo a formatação existente, se fornecido
    if texto2:
        if len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[1]
            if p.runs:
                run = p.runs[0]
                formatar_texto(run, texto2, fonte2, tamanho2)
            else:
                p.text = texto2
                p.font.name = fonte2
                p.font.size = Pt(tamanho2)
                p.alignment = PP_ALIGN.CENTER  # Alinha o texto no centro

        else:
            p = text_frame.add_paragraph()
            p.text = texto2
            p.font.name = fonte2
            p.font.size = Pt(tamanho2)
            p.alignment = PP_ALIGN.CENTER  # Alinha o texto no centro


# Função para atualizar gráficos
def atualizar_grafico(shape, categorias, valores_cativo, valores_livre, fonte="Calibri", tamanho=18):
    chart = shape.chart
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    chart_data.add_series("Cativo", valores_cativo)
    chart_data.add_series("Livre", valores_livre)
    chart.replace_data(chart_data)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(tamanho)  # Usar o tamanho passado como parâmetro
    category_axis.tick_labels.font.name = fonte  # Usar a fonte passada como parâmetro


def atualizar_grafico_com_dados(chart, categorias, *series):
    """Atualiza o gráfico com categorias e séries de dados."""
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    for nome, valores in series:
        chart_data.add_series(nome, valores)
    chart.replace_data(chart_data)


'''
Tarifas ENEL CE:
'''
TE_ENEL_CE = 0.27291
TUSD_ENEL_CE = 0.44929

ICMS = 0.2
PISCOF = 0.0663

tarifa_simpostos = TE_ENEL_CE + TUSD_ENEL_CE

tarifa_cimpostos = round(tarifa_simpostos / ((1 - ICMS) * (1 - (PISCOF))), 4)

ICMS_TE = (0.2 * TE_ENEL_CE) / ((1 - ICMS) * (1 - (PISCOF)))
ICMS_TUSD = (0.2 * TUSD_ENEL_CE) / ((1 - ICMS) * (1 - (PISCOF)))

Tarifa_compensavel_GC = round(tarifa_cimpostos - ICMS_TE - ICMS_TUSD, 4)
Tarifa_compensavel_real = round(tarifa_cimpostos - ICMS_TUSD, 4)

'''
Variáveis:
'''


def powerpoint_edit(infos,buffer):

    desconto = infos['desconto']
    consumokwh = infos['consumo']
    custo_disponibilidade = infos['custo_disponibilidade']
    fidelidade = infos['fidelidade']
    cliente = infos['nome_cliente']

    '''
    Calculos
    '''
    desc = round(desconto / 100, 4)

    tarifa_resultante = round(Tarifa_compensavel_real * (1 - desc), 4)

    consumoRS = round(consumokwh * tarifa_cimpostos, 2)
    consumoANO = round(consumoRS * 12, 2)

    custo_disponibilidadeRS = round(custo_disponibilidade * tarifa_cimpostos, 2)
    goenerRS = (consumokwh - custo_disponibilidade) * tarifa_resultante

    ICMS_RS = round((consumokwh - custo_disponibilidade) * ICMS_TUSD, 2)

    custo_total_GD = round(goenerRS + custo_disponibilidadeRS + ICMS_RS, 2)
    custo_total_GD_ANO = round(custo_total_GD * 12, 2)

    economia_mensal = round(consumoRS - custo_total_GD, 2)
    economia_anual = round(economia_mensal * 12, 2)
    economia_percentual = round((economia_anual / (consumoRS * 12)) * 100, 2)

    # print(f'Tarifa sem impostos: {tarifa_simpostos}')
    # print(f'Tarifa com impostos: {tarifa_cimpostos}')
    # print(f'Tarifa compensável GC: {Tarifa_compensavel_GC}')
    # print(f'Tarifa compensável Atual: {Tarifa_compensavel_real}')
    # print(f'Tarifa resultante (com desconto): {tarifa_resultante}\n')
    #
    # print(f'ICMS: {ICMS_RS}')
    # print(f'Custo mensal sem Goener: {consumoRS}')
    # print(f'Custo mensal com Goener: {custo_total_GD}\n')
    #
    # print(f'Custo anual sem Goener: {consumoANO}')
    # print(f'Custo anual Ccom Goener: {custo_total_GD * 12}\n')
    #
    # print(f'Economia mensal: {economia_mensal}')
    # print(f'Economia anual: {economia_anual}')
    # print(f'Economia percentual: {economia_percentual}')

    economia_ano2 = round(economia_anual * 1.05, 2)
    economia_ano3 = round(economia_ano2 * 1.05, 2)
    economia_ano4 = round(economia_ano3 * 1.05, 2)
    economia_ano5 = round(economia_ano4 * 1.05, 2)
    economia_5anos = economia_anual + economia_ano2 + economia_ano3 + economia_ano4 + economia_ano5

    # Transformar a variável 'economia_anual' em uma string sem o prefixo 'b' e com codificação 'utf-8'
    economia_anual_str = f"{economia_anual:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    economia_5anos_str = f"{economia_5anos:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    economia_ano2_str = f"{economia_ano2:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    economia_ano3_str = f"{economia_ano3:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    economia_ano4_str = f"{economia_ano4:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    economia_ano5_str = f"{economia_ano5:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    economia_mensal_str = f"{economia_mensal:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    consumoANO_str = f"{consumoANO:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    custo_total_GD_ANO_str = f"{custo_total_GD_ANO:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    #print(economia_mensal_str)

    '''
    Atualização da proposta:
    '''
    # Processamento do tempo de fidelidade:
    if fidelidade == 0:
        fidelidade_str = "Sem período de fidelidade"
    else:
        fidelidade_str = f'{fidelidade} meses de fidelidade'

    # Obter a data de hoje no formato de texto
    data = datetime.datetime.now().strftime("%d/%m/%Y")

    # Carregar a apresentação
    ap = Presentation(buffer)

    # Atualizar textos e gráficos
    for slide in ap.slides:
        for shape in slide.shapes:
            if shape.shape_type == 17:  # 17 representa "Text Box"
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("PROPOSTA PARA"):  # OK
                        # print('entrou proposta de energia')
                        texto_atualizado = f"{cliente}"
                        # Limpar o texto existente
                        # shape.text_frame.clear()
                        # Atualizar com o novo texto
                        atualizar_texto_shape_proposta(shape, texto1="Proposta de energia por assinatura:",
                                                       texto2=texto_atualizado, fonte1="Arial", tamanho1=18,
                                                       fonte2="Arial Bold", tamanho2=18)
                    elif text.startswith("XX%"):  # OK
                        # print('entrou % desconto')
                        texto_atualizado = f"{desconto}% DE DESCONTO"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial Bold", tamanho1=16)
                    elif text.startswith("R$ A.BBB,CC"):  # OK
                        # print('entrou economia prevista ano')
                        texto_atualizado = f"R$ {economia_anual_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Open Sans Bold", tamanho1=34)
                    elif text.startswith("R$ X.XXX,XX"):  # OK
                        # print('entrou economia 5 anos')
                        texto_atualizado = f"R$ {economia_5anos_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Open Sans Bold", tamanho1=16)
                    elif text.startswith("R$ BB.CCC"):  # OK
                        # print('entrou CUSTO TOTAL ATUAL')
                        texto_atualizado = f"R$ {consumoANO_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Open Sans Bold", tamanho1=16)
                    elif text.startswith("R$ CC.BBB"):  # OK
                        # print('entrou custo total com Goener')
                        texto_atualizado = f"R$ {custo_total_GD_ANO_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Open Sans Bold", tamanho1=16)
                    elif text.startswith("Sem fidelidade"):  # OK
                        # print('entrou Fidelidade')
                        # shape.text_frame.clear()
                        texto_atualizado = f"{fidelidade_str}, sem investimento, sem instalação ou modificação na rede elétrica atual"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=10)
                    elif text.startswith("00/00"):  # OK
                        # print('entrou data')
                        texto_atualizado = f"{data}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial italics", tamanho1=8)
                    elif text.startswith("R$ 1.aaa,05"):
                        #print('entrou economia ano 1')
                        texto_atualizado = f"R$ {economia_anual_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=5)
                    elif text.startswith("R$ 1.bbb,05"):
                        #print('entrou economia ano 2')
                        texto_atualizado = f"R$ {economia_ano2_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=5)
                    elif text.startswith("R$ 1.ccc,05"):
                        #print('entrou economia ano 3')
                        texto_atualizado = f"R$ {economia_ano3_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=5)
                    elif text.startswith("R$ 1.ddd,05"):
                        #print('entrou economia ano 4')
                        texto_atualizado = f"R$ {economia_ano4_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=5)
                    elif text.startswith("R$ 1.eee,05"):
                        #print('entrou economia ano 5')
                        texto_atualizado = f"R$ {economia_ano5_str}"
                        atualizar_texto_shape(shape, texto1=texto_atualizado, fonte1="Arial", tamanho1=5)

    # Salvar as alterações em novo buffer
    output_buffer = io.BytesIO()
    ap.save(output_buffer)
    output_buffer.seek(0)  # Voltar o cursor para o início do buffer

    return output_buffer
    # pptx_filename = f"{cliente} - Proposta de Energia por assinatura.pptx"
    # ap.save(pptx_filename)

    # # Verificar se o arquivo PPTX foi salvo corretamente
    # if os.path.exists(pptx_filename):
    #     # Converter o arquivo PPTX em PDF
    #     pdf_filename = f"{cliente} - Proposta de Energia por assinatura.pdf"
    #     convert_ppt_to_pdf(pptx_filename, pdf_filename)
    # else:
    #     print(f"Erro: O arquivo {pptx_filename} não foi encontrado.")